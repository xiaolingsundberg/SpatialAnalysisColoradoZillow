from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


PROJECT_DIR = Path("/Users/xiaolingsundberg/Desktop/STA6709/SemesterProject")
WRITEUP_DIR = PROJECT_DIR / "WRITE UP"
ASSET_DIR = WRITEUP_DIR / "revised_report_assets"
OUT = WRITEUP_DIR / "STA6709_Colorado_Zillow_Spatial_Analysis_Presentation.pptx"

BLUE = RGBColor(31, 78, 121)
DARK = RGBColor(31, 31, 31)
GRAY = RGBColor(96, 96, 96)
LIGHT = RGBColor(242, 246, 249)
ACCENT = RGBColor(48, 160, 120)
ORANGE = RGBColor(242, 145, 44)


def add_textbox(slide, text, left, top, width, height, font_size=22, bold=False,
                color=DARK, align=PP_ALIGN.LEFT, font="Aptos"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, title, Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.55),
                font_size=28, bold=True, color=BLUE)
    if subtitle:
        add_textbox(slide, subtitle, Inches(0.58), Inches(0.92), Inches(11.8), Inches(0.35),
                    font_size=12, color=GRAY)
    line = slide.shapes.add_shape(1, Inches(0.58), Inches(1.22), Inches(12.05), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(210, 220, 230)
    line.line.color.rgb = RGBColor(210, 220, 230)


def add_footer(slide, n):
    add_textbox(slide, f"STA6709 Spatial Statistics | {n}", Inches(11.35), Inches(7.08),
                Inches(1.3), Inches(0.2), font_size=8, color=GRAY, align=PP_ALIGN.RIGHT)


def add_bullets(slide, bullets, left, top, width, height, font_size=18):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(10)
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
    return box


def add_metric(slide, label, value, x, y, w=Inches(2.3), h=Inches(0.92), color=BLUE):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = RGBColor(210, 220, 230)
    add_textbox(slide, value, x + Inches(0.10), y + Inches(0.10), w - Inches(0.20), Inches(0.35),
                font_size=20, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, x + Inches(0.10), y + Inches(0.52), w - Inches(0.20), Inches(0.25),
                font_size=9, color=GRAY, align=PP_ALIGN.CENTER)


def add_picture(slide, filename, left, top, width=None, height=None):
    path = ASSET_DIR / filename
    return slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def add_table(slide, data, left, top, width, height, header_fill=RGBColor(216, 224, 232),
              font_size=10):
    rows, cols = len(data), len(data[0])
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for r, row in enumerate(data):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Aptos"
                p.font.size = Pt(font_size)
                p.font.color.rgb = DARK
                if r == 0:
                    p.font.bold = True
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
    return table


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# 1. Title
slide = prs.slides.add_slide(blank)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(250, 252, 253)
add_textbox(slide, "Spatial Analysis of Colorado Zillow Housing Listings",
            Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.9),
            font_size=34, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Mapping, regression, nearest-neighbor analysis, clustering, and interactive Leaflet visualization",
            Inches(1.2), Inches(2.25), Inches(10.9), Inches(0.5),
            font_size=17, color=GRAY, align=PP_ALIGN.CENTER)
add_metric(slide, "cleaned listings", "123", Inches(2.2), Inches(3.55), color=BLUE)
add_metric(slide, "cities represented", "44", Inches(5.5), Inches(3.55), color=ACCENT)
add_metric(slide, "ZIP codes", "82", Inches(8.8), Inches(3.55), color=ORANGE)
add_textbox(slide, "Xiaoling Sundberg | STA6709 Spatial Statistics",
            Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.35),
            font_size=13, color=GRAY, align=PP_ALIGN.CENTER)
add_footer(slide, 1)

# 2. Objective
slide = prs.slides.add_slide(blank)
add_title(slide, "Project Objective", "Use spatial methods to understand where listings are located and how prices vary.")
add_bullets(slide, [
    "Treat Zillow listings as spatial points using latitude and longitude.",
    "Map listings within Colorado state and county polygon boundaries.",
    "Use regression to evaluate how price relates to housing characteristics and location.",
    "Use nearest-neighbor and cluster analysis to assess whether listings are random or geographically concentrated.",
    "Create an interactive Leaflet map for exploring individual listings."
], Inches(0.8), Inches(1.75), Inches(6.1), Inches(4.5), font_size=19)
add_picture(slide, "figure1_colorado_boundary_house_points.png", Inches(7.15), Inches(1.55), width=Inches(5.3))
add_footer(slide, 2)

# 3. Dataset
slide = prs.slides.add_slide(blank)
add_title(slide, "Dataset and Preparation", "The data were cleaned before mapping, modeling, and spatial analysis.")
add_table(slide, [
    ["Metric", "Value"],
    ["Number of listings", "123"],
    ["Number of cities", "44"],
    ["Number of ZIP codes", "82"],
    ["Minimum price", "$49,999"],
    ["Median price", "$635,000"],
    ["Mean price", "$3,972,322"],
    ["Maximum price", "$300,000,000"],
    ["Median area square feet", "2,540"],
    ["Median price per square foot", "$254"],
], Inches(0.78), Inches(1.55), Inches(5.55), Inches(4.85), font_size=12)
add_bullets(slide, [
    "Duplicate listing IDs were removed.",
    "Records with missing price, coordinates, beds, baths, or square footage were excluded.",
    "Price per square foot and log price were created for analysis.",
    "Coordinates were checked to keep listings inside a reasonable Colorado range."
], Inches(6.75), Inches(1.75), Inches(5.75), Inches(3.8), font_size=18)
add_footer(slide, 3)

# 4. Price distribution
slide = prs.slides.add_slide(blank)
add_title(slide, "Exploratory Summary: Price Is Highly Skewed", "A small number of luxury listings strongly affects the mean.")
add_picture(slide, "figure4_price_histogram.png", Inches(0.75), Inches(1.45), width=Inches(5.75))
add_picture(slide, "figure5_price_by_area.png", Inches(6.85), Inches(1.45), width=Inches(5.75))
add_textbox(slide, "Most listings fall in moderate price ranges, but a few very high-value homes create a long right tail. This is why log listing price was used in the regression model.",
            Inches(1.0), Inches(6.4), Inches(11.4), Inches(0.45), font_size=16, color=DARK, align=PP_ALIGN.CENTER)
add_footer(slide, 4)

# 5. Static maps
slide = prs.slides.add_slide(blank)
add_title(slide, "Listings Are Concentrated Along the Front Range", "The map shows a strong geographic concentration rather than an even statewide pattern.")
add_picture(slide, "figure2_county_boundary_house_points.png", Inches(0.9), Inches(1.45), width=Inches(6.25))
add_bullets(slide, [
    "Listings are concentrated near Denver, Colorado Springs, Boulder, and surrounding communities.",
    "County polygons provide geographic context for the listing points.",
    "The spatial pattern suggests housing data should be interpreted geographically, not only numerically."
], Inches(7.35), Inches(1.85), Inches(5.15), Inches(3.35), font_size=19)
add_footer(slide, 5)

# 6. Home type
slide = prs.slides.add_slide(blank)
add_title(slide, "Home Type Pattern", "Single-family homes dominate the sample, so smaller groups should be interpreted cautiously.")
add_picture(slide, "figure3_home_type_map.png", Inches(0.8), Inches(1.45), width=Inches(6.3))
add_table(slide, [
    ["Home Type", "Listings", "Median Price", "Median Area"],
    ["Single-family", "113", "$650,000", "2,583"],
    ["Condo", "5", "$375,000", "1,201"],
    ["Townhouse", "5", "$585,000", "1,780"],
], Inches(7.45), Inches(1.75), Inches(4.7), Inches(1.55), font_size=12)
add_bullets(slide, [
    "The transparent single-family layer keeps condos and townhouses visible.",
    "Small condo and townhouse counts mean their summary statistics are less stable.",
    "Home type was included as a predictor in the regression model."
], Inches(7.45), Inches(3.65), Inches(4.7), Inches(2.3), font_size=17)
add_footer(slide, 6)

# 7. Regression
slide = prs.slides.add_slide(blank)
add_title(slide, "Regression Analysis", "Log listing price was modeled using property characteristics and location.")
add_table(slide, [
    ["Variable", "Estimate", "p-value"],
    ["Beds", "-0.020", "0.6751"],
    ["Baths", "0.111", "0.0432"],
    ["Area Sq Ft", "0.000195", "<0.001"],
    ["Single Family", "0.442", "0.0486"],
    ["Townhouse", "0.435", "0.1498"],
    ["Days on Zillow", "-0.00477", "0.0001"],
    ["Latitude", "0.163", "0.0236"],
    ["Longitude", "-0.437", "<0.001"],
], Inches(0.75), Inches(1.45), Inches(5.7), Inches(4.65), font_size=11)
add_bullets(slide, [
    "Square footage, bathrooms, days on Zillow, latitude, and longitude were meaningful predictors.",
    "The location coefficients show that geography remained important even after adding property characteristics.",
    "Using log price helped reduce the influence of extreme listing prices."
], Inches(6.95), Inches(1.8), Inches(5.55), Inches(3.25), font_size=19)
add_footer(slide, 7)

# 8. Residual map
slide = prs.slides.add_slide(blank)
add_title(slide, "Regression Residuals by Location", "Residual mapping shows where the model overpredicts or underpredicts price geographically.")
add_picture(slide, "figure6_regression_residual_map.png", Inches(0.9), Inches(1.45), width=Inches(6.25))
add_bullets(slide, [
    "Positive residuals mean actual price was higher than predicted.",
    "Negative residuals mean actual price was lower than predicted.",
    "Large residuals may reflect factors not in the model, such as views, neighborhood quality, school zones, or luxury amenities."
], Inches(7.35), Inches(1.85), Inches(5.1), Inches(3.7), font_size=18)
add_footer(slide, 8)

# 9. Nearest-neighbor
slide = prs.slides.add_slide(blank)
add_title(slide, "Nearest-Neighbor Analysis", "Listings are closer together than expected under complete spatial randomness.")
add_table(slide, [
    ["Metric", "Value"],
    ["Listings", "123"],
    ["Study area", "81,832.085 sq km"],
    ["Observed mean NN distance", "8.356 km"],
    ["Expected mean NN distance", "12.897 km"],
    ["Nearest-neighbor index", "0.648"],
    ["Interpretation", "Clustered"],
], Inches(0.75), Inches(1.55), Inches(4.75), Inches(3.2), font_size=12)
add_picture(slide, "figure7_nearest_neighbor_histogram.png", Inches(5.85), Inches(1.35), width=Inches(6.6))
add_textbox(slide, "Because the index is less than 1, the listings are spatially clustered rather than randomly distributed.",
            Inches(0.95), Inches(5.15), Inches(4.25), Inches(0.8), font_size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_footer(slide, 9)

# 10. Spatial clusters
slide = prs.slides.add_slide(blank)
add_title(slide, "Spatial Cluster Analysis", "Coordinate-based clusters summarize regional listing groups.")
add_picture(slide, "figure8_spatial_clusters.png", Inches(0.75), Inches(1.42), width=Inches(5.95))
add_table(slide, [
    ["Cluster", "Count", "Median Price", "Approx. Location"],
    ["I", "29", "$460,000", "Southern CO / Colorado Springs"],
    ["II", "3", "$75,000,000", "Mountain / luxury listings"],
    ["III", "72", "$662,500", "Denver & central Front Range"],
    ["IV", "19", "$650,000", "Northern Front Range / Boulder"],
], Inches(6.95), Inches(1.65), Inches(5.6), Inches(2.4), font_size=10)
add_bullets(slide, [
    "The largest cluster is around Denver and the central Front Range.",
    "The small luxury cluster has very high prices but only three listings.",
    "Clusters help connect price patterns to regional geography."
], Inches(7.15), Inches(4.45), Inches(5.05), Inches(1.8), font_size=17)
add_footer(slide, 10)

# 11. Interactive map
slide = prs.slides.add_slide(blank)
add_title(slide, "Interactive Leaflet Map", "The final HTML map allows viewers to inspect individual listings.")
add_bullets(slide, [
    "Shows Colorado boundary, county polygons, and Zillow listing points.",
    "Popups include address, city, actual price, predicted price, residual, home type, beds, baths, area, and cluster.",
    "Submitted separately as a zipped HTML map so it can be opened in a web browser."
], Inches(0.85), Inches(1.65), Inches(5.95), Inches(3.7), font_size=19)
shape = slide.shapes.add_shape(1, Inches(7.35), Inches(1.75), Inches(4.55), Inches(2.65))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT
shape.line.color.rgb = RGBColor(205, 218, 230)
add_textbox(slide, "Submitted file", Inches(7.65), Inches(2.05), Inches(3.95), Inches(0.35),
            font_size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_textbox(slide, "colorado_zillow_interactive_map_SUBMIT.zip", Inches(7.65), Inches(2.55),
            Inches(3.95), Inches(0.65), font_size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_textbox(slide, "Unzip, then open colorado_zillow_complete_interactive_map.html", Inches(7.65), Inches(3.35),
            Inches(3.95), Inches(0.55), font_size=12, color=GRAY, align=PP_ALIGN.CENTER)
add_footer(slide, 11)

# 12. Conclusion
slide = prs.slides.add_slide(blank)
add_title(slide, "Conclusion", "Spatial analysis added insight that tables alone could not show.")
add_metric(slide, "nearest-neighbor index", "0.648", Inches(0.9), Inches(1.55), color=BLUE)
add_metric(slide, "largest cluster", "72", Inches(3.55), Inches(1.55), color=ACCENT)
add_metric(slide, "median listing price", "$635K", Inches(6.2), Inches(1.55), color=ORANGE)
add_metric(slide, "mean listing price", "$3.97M", Inches(8.85), Inches(1.55), color=BLUE)
add_bullets(slide, [
    "Colorado Zillow listings in this sample are clustered, especially along the Front Range.",
    "Location remains important after accounting for home characteristics in regression.",
    "Residual and cluster maps help identify geographic patterns that are not obvious from tables.",
    "The project demonstrates how mapping, modeling, and interactive visualization can work together in spatial housing analysis."
], Inches(1.05), Inches(3.05), Inches(11.1), Inches(3.0), font_size=20)
add_footer(slide, 12)

prs.save(OUT)
print(OUT)
